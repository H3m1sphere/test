# 書き換えうる変数 #
temprate_pptx_address = './resources/pp_template.pptx'
layout_name = 'pypptx_layout'

# 略語集 #
# prs: presentation
# ph: placeholder


from pptx import Presentation
from pptx.util import Cm
import matplotlib.pyplot as plt
import pandas as pd
import os



def search_layout(prs, layout_name):
    layout = next((layout for layout in prs.slide_layouts if layout.name == layout_name), None)
    if layout == None:
        print(f"No items matching '{layout_name}' were found.")
    return layout

def search_layoutplaceholder(layout, text):
    layoutplaceholder = next((ph for ph in layout.placeholders if ph.has_text_frame and ph.text == text), None)
    if layoutplaceholder == None:
        print(f"No items matching '{text}' were found.")
    return layoutplaceholder

def search_slideplaceholder(slide, text):
    # レイアウト上のphの位置とサイズを取得し，スライド上の位置が一致するphを返す．
    # 迂遠だが，スライド上のphを文字列で指定したいが，その手段が分からないのでこうする．
    layout_ph = search_layoutplaceholder(slide.slide_layout, text)
    if layout_ph is None:
        print(f"No layout placeholder matching '{text}' was found.")
        return None

    for shape in slide.shapes:
        if (shape.left, shape.top, shape.width, shape.height) == (layout_ph.left, layout_ph.top, layout_ph.width, layout_ph.height):
            return shape
    print(f"No items matching {slide}'s shapes and '{text}' were found.")

def show_placeholders_in_slide(slide):
    print(f'{slide.name} has {len(slide.placeholders)} placeholders.')
    for ph in slide.placeholders:
        print(f"ID: {ph.placeholder_format.idx}, Type: {ph.placeholder_format.type}, Text: {ph.text}")

def show_shapes_in_slide(slide):
    print(f'{slide.name} has {len(slide.shapes)} shapes.')
    for shape in slide.shapes:
        print(f"ID: {shape.shape_id}, Type: {shape.shape_type}, Name: {shape.name}, Text: {shape.text if shape.has_text_frame else ''}")

def make_graph():
    # グラフのデータ
    x = [1, 2, 3, 4, 5]
    y = [1, 4, 9, 16, 25]

    # グラフを作成
    plt.figure(figsize=(6, 4))
    plt.plot(x, y, marker='o')
    plt.title("Sample Plot")
    plt.xlabel("X axis")
    plt.ylabel("Y axis")


    # グラフを画像ファイルとして保存
    plt.savefig('./temp/temp_graph.png')
    plt.close()

def make_table():
    data = {'Name': ['John', 'Anna', 'Peter', 'Linda'],
        'Occupation': ['Engineer', 'Doctor', 'Architect', 'Teacher'],
        'Age': [28, 34, 45, 32]}
    df = pd.DataFrame(data) 
    return(df)

# MAIN #

# テンプレートとなるpptxを読み込み，レイアウトを指定して，スライドを作成する．
prs = Presentation(temprate_pptx_address)
slide_layout = search_layout(prs, layout_name) # 「レイアウト名の変更」を参照すること
slide = prs.slides.add_slide(slide_layout)
# shapes確認用テスト関数
# show_shapes_in_slide(ppteset_slide)
# show_shapes_in_slide(slide)



# 作成したスライドのPlaceholderを取得．Placeholder上のtextで探索している．
title_ph = slide.shapes.title
graph_ph = search_slideplaceholder(slide, "graph")
table_ph = search_slideplaceholder(slide, "table")



# グラフを作成し，画像ファイルとして保存，画像ファイルをスライドに挿入．
make_graph()
graph_ph.insert_picture('./temp/temp_graph.png')
os.remove('./temp/temp_graph.png')



df = make_table()
# 表プレースホルダーにデータフレームの内容を挿入
graphic_frame = table_ph.insert_table(rows=df.shape[0]+1, cols=df.shape[1])
table = graphic_frame.table  # GraphicFrameからTableオブジェクトを取得

# ヘッダー行の追加
for col_index, col_name in enumerate(df.columns):
    table.cell(0, col_index).text = col_name

# データ行の追加
for row_index, row in df.iterrows():
    for col_index, item in enumerate(row):
        table.cell(row_index + 1, col_index).text = str(item)

# 作成したスライドを保存
prs.save('./output/test.pptx')